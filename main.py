from pptx import Presentation
from pptx.util import Inches
import qrcode

def replace_shape_with_qrcode(slide, shape_index, qr_code_data):
    # Get the shape to replace
    shape = slide.shapes[shape_index]

    # Remove existing shape
    sp = slide.shapes._spTree.remove(shape._element)

    # Create QR code image
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_L,
        box_size=10,
        border=4,
    )
    qr.add_data(qr_code_data)
    qr.make(fit=True)

    # Create QR code image and insert it into the slide
    img = qr.make_image(fill_color="black", back_color="white")
    image_path = "temp_qr_code.png"
    img.save(image_path)
    slide.shapes.add_picture(image_path, shape.left, shape.top, width=shape.width, height=shape.height)


def replace_name(slide, name,font): 
    for i in range(len(slide.shapes)):
        if slide.shapes[i].has_text_frame:
            if slide.shapes[i].text == "[CLASS + NUMBER]":
                text_frame = slide.shapes[i].text_frame
                text_frame.text = text_frame.text.replace("[CLASS + NUMBER]", name)
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = font
                return "Replacement Sucesfull"
    raise Exception("No name found")
    
def printShapes(slide):
    for i, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            print(f"Shape {i+1} (Text): {shape.text_frame.text}")
        else:
            print(f"Shape {i+1} (Non-Text)")

def main():
    # Load the PowerPoint presentation
    presentation = Presentation("your_template.pptx")
    Intro_Slide = presentation.slides[0] 
    Indivual_Slide = presentation.slides[1]
    Group_Slide = presentation.slides[2]

    # Specify the slide index and shape index you want to replace with QR 
    slide_index = 0
    shape_index = 0

    # QR code data (replace with actual data)
    qr_code_data = "https://www.example.com"

    # Replace the name of the class
    replace_name(Intro_Slide, "CSE 100", "Times New Roman") 

    # Replace the shape with QR code
    replace_shape_with_qrcode(Indivual_Slide,4,qr_code_data)
    replace_shape_with_qrcode(Group_Slide,3,qr_code_data)

    # Save the updated presentation 
    presentation.save("output_presentation.pptx")

   
    # printShapes(Indivual_Slide)

if __name__ == "__main__":
    main()
