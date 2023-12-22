from pptx import Presentation
import os
import json

# Load JSON data from file
with open('test.json', 'r') as file:
    data = json.load(file)
    print(data)

# Create a new PowerPoint presentation
presentation = Presentation()

# Iterate over the JSON data and add slides to the presentation
for item in data:
    layout_name = item['layout']
    
    # Check if the layout exists in the new presentation
    existing_layouts = [layout.name for layout in presentation.slide_layouts]
    if layout_name not in existing_layouts:
        print(f"Error: Slide layout '{layout_name}' not found. Skipping slide.")
        continue

    new_layout = presentation.slide_layouts.get_by_name(layout_name)
    slide = presentation.slides.add_slide(new_layout)

    for shape_data in item['shapes']:
        shape_type = shape_data['type']

        # Check if shape type is valid (e.g., 14 is for title)
        if shape_type != 14:
            print(f"Error: Invalid shape type '{shape_type}'. Skipping shape.")
            continue

        # Add title to the title placeholder
        title_shape = slide.shapes.title
        title_shape.text = shape_data['text']

# Save the presentation as a PPTX file
presentation.save('output.pptx')